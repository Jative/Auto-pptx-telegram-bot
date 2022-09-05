import requests
import telebot
import glob
from PIL import Image, ImageSequence
from bs4 import BeautifulSoup
from pptx import Presentation
from os import remove, path
from time import sleep


token = "1968458319:AAEod7XMwEIprViQ6BFn09gG9oRR7bjW7zk"
parsed_text = [[]]
parsed_titles = []
title_black_list = ["Классификация", "Аксессуары", "Примечания", "Литература", "Литература по теме", "Ссылки",
                    "Галерея", "Специализирующиеся производители", "См. также", "Награды", "Города-побратимы"]
WIKI_URL = "https://ru.wikipedia.org/wiki/"
max_text_len = 1100
max_links = 600
mp_error = False
actual_content = {}
number_of_slides = 0


def wiki_parser(url):
    global parsed_text, parsed_titles

    q = requests.get(url)
    result = q.content
    soup = BeautifulSoup(result, "lxml")

    all_info = soup.find(class_="mw-parser-output")

    title_count = 0

    for child in all_info.children:
        if child.name == "p" and not isinstance(child, str):
            try:
                text = child.text
                if text == "#invoke:Navbox\n" or "См. также:" in text:
                    continue
                for i in range(1, max_links):
                    text = text.replace(f"[{i}]", "").replace(f"[комм. {i}]", "")
                if text[-1] == ":":
                    text = text[:-1] + "."
                parsed_text[title_count].append(text)
            except:
                continue
        elif child.name == "h2" and not isinstance(child, str):
            text = child.text.replace("[править | править код]", "").replace("[значимость факта?]", "")
            parsed_titles.append(text)
            title_count += 1
            parsed_text.append([])
        elif child.name == "div" or child.name == "table":
            try:
                content_class = child.get("class")
                if title_count:
                    actual_title = parsed_titles[-1]
                    if actual_title not in actual_content and "thumb" in content_class or "infobox" in content_class:
                        content_href = "https:"+child.find("img").get("src")
                        content_type = content_href[-3:]
                        r = requests.get(content_href)
                        content_name = f"{actual_title}.{content_type}"
                        with open(f"data/{content_name}", "bw") as f:
                            f.write(r.content)
                        actual_content[actual_title] = content_name
                else:
                    actual_title = 0
                    if actual_title not in actual_content and "thumb" in content_class or "infobox" in content_class:
                        content_href = "https:"+child.find("img").get("src")
                        content_type = content_href[-3:]
                        r = requests.get(content_href)
                        content_name = f"{actual_title}.{content_type}"
                        with open(f"data/{content_name}", "bw") as f:
                            f.write(r.content)
                        actual_content[actual_title] = content_name
            except:
                continue

    for text in parsed_text:
        for p in text:
            if p != "\n":
                break
        else:
            local_ind = parsed_text.index(text)
            if local_ind != 0:
                del parsed_text[local_ind]
                del parsed_titles[local_ind-1]
    while True:
        if [] in parsed_text:
            ind = parsed_text.index([])
            del parsed_text[ind]
            del parsed_titles[ind-1]
        else:
            break


def clean_trash():
    for title in parsed_titles:
        ind = parsed_titles.index(title)+1
        if "Известные" in title:
            del parsed_text[ind]
            del parsed_titles[ind-1]
        if title in title_black_list:
            all_text = ""
            for text in parsed_text[ind]:
                all_text += text
            if len(all_text) < 50:
                del parsed_text[ind]
                del parsed_titles[ind-1]


def thumbnails(frames):
    for frame in frames:
        thumbnail = frame.copy()
        thumbnail.thumbnail(size, Image.ANTIALIAS)
        yield thumbnail


def content_resize():
    for key in actual_content.keys():
        try:
            correct_dir = f"data/{actual_content[key]}"
            content = Image.open(correct_dir)
            x, y = content.size

            if y / x != 10 / 6.7:
                if x > y:
                    new_x = x
                    new_y = int(x * 1.49)

                elif x < y or x == y:
                    new_x = int(y / 1.49)
                    new_y = y

                if actual_content[key][-3:] == "gif":
                    frames = ImageSequence.Iterator(content)
                    frames = thumbnails(frames)
                    gif_content = next(frames)
                    gif_content.info = content.info
                    gif_content.save(correct_dir, save_all=True, append_images=list(frames))

                else:
                    result = content.crop((0, 0, x, y)).resize((new_x, new_y))
                    result.save(correct_dir)

        except:
            continue



def delete_slide(index):
    global prs
    xml_slides = prs.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def presentation_maker():
    global actual_content, prs, number_of_slides

    clean_trash()
    content_resize()

    try:
        prs = Presentation("sample.pptx")
    except:
        prs = Presentation()

    first_slide_layout = prs.slide_layouts[0]
    first_slide = prs.slides.add_slide(first_slide_layout)
    first_count = 0
    slides_with_error = []

    for shape in first_slide.shapes:
        if not shape.has_text_frame:
            continue
        if not first_count:
            shape.text = theme
            first_count += 1
        else:
            shape.text = f"Подготовил {name}"

    second_slide_layout = prs.slide_layouts[1] if 0 not in actual_content else prs.slide_layouts[8]
    second_slide = prs.slides.add_slide(second_slide_layout)
    if second_slide_layout == prs.slide_layouts[8]:
        try:
            placeholder = second_slide.placeholders[1]
            placeholder.insert_picture(f"data/{actual_content[0]}")
        except:
            slides_with_error.append(1)
            second_slide_layout = prs.slide_layouts[1]
            second_slide = prs.slides.add_slide(second_slide_layout)

    second_count = 0
    for shape in second_slide.shapes:
        if not shape.has_text_frame:
            continue
        elif not second_count:
            shape.text = theme
            second_count += 1
        else:
            text = "\n".join(parsed_text[0])
            if len(text) > max_text_len:
                symbol_counter = -1
                for symbol in text:
                    symbol_counter += 1
                    if symbol == ".":
                        dot_ind = symbol_counter
                    if symbol_counter > max_text_len:
                        break
                text = text[:dot_ind+1]
            shape.text = text

    for title in parsed_titles:
        slide_layout = prs.slide_layouts[1] if title not in actual_content else prs.slide_layouts[8]
        title_ind = parsed_titles.index(title)
        slide = prs.slides.add_slide(slide_layout)
        if slide_layout == prs.slide_layouts[8]:
            try:
                placeholder = slide.placeholders[1]
                placeholder.insert_picture(f"data/{actual_content[title]}")
            except:
                slides_with_error.append(title_ind+2)
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

        local_count = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            elif not local_count:
                shape.text = title
                local_count += 1
            else:
                text = "\n".join(parsed_text[title_ind+1])
                if len(text) > max_text_len:
                    symbol_counter = -1
                    for symbol in text:
                        symbol_counter += 1
                        if symbol == ".":
                            dot_ind = symbol_counter
                        if symbol_counter > max_text_len:
                            break
                    text = text[:dot_ind+1]
                shape.text = text
                local_count = 0

    for slide_with_error_ind in slides_with_error:
        delete_slide(slide_with_error_ind)

    prs.save(f'{theme}.pptx')
    number_of_slides = len(prs.slides)
    print("Количество слайдов:", number_of_slides)



def make_presentation():
    global mp_error, actual_content

    try:
        #Забираю данные с Википедии
        wiki_parser(WIKI_URL+theme.replace(" ", "_"))
        print("Парсинг завершён")

        #Создаю презентацию
        presentation_maker()
    except:
        mp_error = True

    actual_content = {}
    files = glob.glob("data/*")
    for file in files:
        remove(file)


def telegram_bot(token):
    global bot
    print("Бот запущен")
    bot = telebot.TeleBot(token, threaded=False)

    @bot.message_handler(commands=["start"])
    def start_message(message):
        try:
            bot.send_message(message.chat.id, "Укажите в первой строчке тему, а во второй имя автора")
        except:
            print("Ошибка отправки(Стартовое сообщение)")

    @bot.message_handler(content_types=["text"])
    def send_text(message):
        global theme, name, parsed_text, parsed_titles, mp_error

        if message.text.count("\n") == 1:
            theme, name = message.text.split("\n")
            print(f"\n{theme}  //////  {name}")
            make_presentation()
            if mp_error:
                try:
                    bot.send_message(message.chat.id, "Извините, у нас технические шоколадки. Пожалуйста, попробуйте позже")
                except:
                    print("Ошибка отправки(Оповещение о несозданной презентации)")
                print("Ошибка создания презентации")
                mp_error = False
            else:
                file = open(f"{theme}.pptx", "rb")
                try:
                    bot.send_document(message.chat.id, file)
                    bot.send_message(message.chat.id, f"Количество слайдов: {number_of_slides}")
                    print("Успешно")
                except:
                    print("Ошибка отправки(Презентация)")
                file.close()
                remove(f"{theme}.pptx")
            parsed_text, parsed_titles = [[]], []
        else:
            try:
                bot.send_message(message.chat.id, "Кажется, что-то пошло не так. Пожалуйста, проверьте форму указанных данных")
            except:
                print("Ошибка отправки(Оповещение о неправильной форме указанных данных)")

    bot.polling()

if __name__ == "__main__":
    telegram_bot(token)
